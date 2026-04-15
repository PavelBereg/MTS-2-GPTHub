"""
PPTX Generator — эндпоинт для генерации профессиональных презентаций.

Маршруты:
  POST /v1/presentations/generate   — создаёт PPTX-файл по теме
  GET  /v1/presentations/download/{file_id} — скачивает готовый файл
  GET  /v1/presentations/list/{user_id}     — список презентаций пользователя

Дизайн:
  - Корпоративные цвета МТС: красный #E30613, тёмный #1C1C1C, белый #FFFFFF
  - Слайды: титульный + контентные (до 10) + заключительный
  - Данные для слайдов генерирует LLM (llama-3.3-70b-instruct)

Интеграция в чат:
  Если classify_intent() → 'presentation', routes.py запускает интерактивный
  presentation-agent. Файл создаётся только после подтверждения структуры.
"""

from __future__ import annotations

import asyncio
import json
import logging
import os
import re
import time
import uuid
from pathlib import Path
from typing import Any, Optional

import httpx
from fastapi import APIRouter
from fastapi.responses import FileResponse, JSONResponse
from pydantic import BaseModel, Field

from app.core.config import get_settings

logger = logging.getLogger("mts.pptx")

router = APIRouter()

# Временная директория для хранения PPTX-файлов
PPTX_DIR = Path("/tmp/mts_presentations")
PPTX_DIR.mkdir(parents=True, exist_ok=True)

# In-memory реестр файлов: file_id → Path
_file_registry: dict[str, dict] = {}


def _shorten_text(text: Any, limit: int) -> str:
    """Сжимает текст до безопасной длины для PPTX-блоков."""
    value = re.sub(r"\s+", " ", str(text or "")).strip()
    if len(value) <= limit:
        return value
    cut = value[: max(0, limit - 1)].rstrip()
    if " " in cut:
        cut = cut.rsplit(" ", 1)[0]
    return cut.rstrip(".,;:") + "…"


def _compact_title(text: str, limit: int = 72) -> str:
    """Убирает мусорные вводные и делает заголовок пригодным для титульника."""
    value = re.sub(r"\s+", " ", str(text or "")).strip(" .,:;")
    value = re.sub(
        r"^(сделай|создай|подготовь|сгенерируй|собери)\s+",
        "",
        value,
        flags=re.IGNORECASE,
    )
    value = re.sub(
        r"\b(презентацию|презентация|слайды|pptx|powerpoint)\b",
        "",
        value,
        flags=re.IGNORECASE,
    )
    value = re.sub(r"\s+", " ", value).strip(" .,:;")
    return _shorten_text(value or "Презентация", limit)


def _compact_subtitle(audience: str, goal: str, deck_subtitle: str = "") -> str:
    if deck_subtitle:
        return _shorten_text(deck_subtitle, 92)
    parts = []
    if audience:
        parts.append(f"Аудитория: {_shorten_text(audience, 42)}")
    if goal:
        parts.append(f"Цель: {_shorten_text(goal, 44)}")
    return _shorten_text(" | ".join(parts), 96)


# ──────────────────────────────────────────────────────────────────
# Pydantic-модели
# ──────────────────────────────────────────────────────────────────


class PresentationRequest(BaseModel):
    topic: str = Field(..., description="Тема презентации", min_length=3, max_length=500)
    slides_count: int = Field(default=5, ge=2, le=15, description="Количество слайдов (2-15)")
    style: str = Field(
        default="modern_light",
        description="Стиль: mts_corporate | strict_corporate | modern_light | editorial | tech | creative | minimal | dark",
    )
    language: str = Field(default="ru", description="Язык слайдов: ru | en")
    user_id: str = Field(default="anonymous", description="ID пользователя")


class OutlineSlideRequest(BaseModel):
    """Слайд, уже согласованный интерактивным PPTX-агентом."""

    title: str
    purpose: str = ""
    bullets: list[str] = Field(default_factory=list)
    visual_hint: str = ""
    speaker_notes: str = ""
    layout: str = ""
    highlight: str = ""
    body: str = ""
    evidence: list[str] = Field(default_factory=list)
    takeaway: str = ""
    icon: str = ""


class PresentationFromOutlineRequest(BaseModel):
    """
    Генерация файла из готового брифа и структуры.

    Этот контракт нужен интерактивному агенту: LLM сначала согласует структуру
    с пользователем, и только потом python-pptx получает финальные данные.
    """

    brief: dict[str, Any] = Field(default_factory=dict)
    outline: list[OutlineSlideRequest] = Field(default_factory=list)
    style: str = Field(
        default="modern_light",
        description="mts_corporate | strict_corporate | modern_light | editorial | tech | creative | minimal | dark",
    )
    language: str = Field(default="ru")
    user_id: str = Field(default="anonymous")


class SlideData(BaseModel):
    title: str
    bullets: list[str]
    notes: str = ""
    purpose: str = ""
    visual_hint: str = ""
    layout: str = ""
    highlight: str = ""
    body: str = ""
    evidence: list[str] = Field(default_factory=list)
    takeaway: str = ""
    icon: str = ""


class PresentationData(BaseModel):
    presentation_title: str
    subtitle: str = ""
    slides: list[SlideData]
    conclusion: str


# ──────────────────────────────────────────────────────────────────
# Генерация структуры через LLM
# ──────────────────────────────────────────────────────────────────

_STRUCTURE_PROMPT = """\
Создай структуру профессиональной презентации на тему: «{topic}».

Требования:
- Количество контентных слайдов: {slides_count} (не считая титульного и заключительного)
- Язык: {language}
- Каждый слайд должен иметь свою роль: проблема, контекст, механизм, сравнение, решение, вывод
- Не делай набор мелких общих фактов. Строй историю, аргументы и причинно-следственные связи
- Каждый слайд: заголовок + сильная крупная мысль + 2-4 коротких буллита
- Используй разные layout: hero, split, cards, process, comparison, quote, takeaway
- Заключение: 2-3 ключевых вывода

Верни ТОЛЬКО валидный JSON без лишнего текста:
{{
  "presentation_title": "...",
  "subtitle": "...",
  "slides": [
    {{
      "title": "Заголовок слайда",
      "layout": "hero|split|cards|process|comparison|quote|takeaway",
      "highlight": "главная крупная фраза слайда",
      "body": "2-3 предложения с обработанным содержанием слайда",
      "evidence": ["конкретный факт или деталь", "ещё один факт"],
      "bullets": ["Буллит 1", "Буллит 2", "Буллит 3"],
      "purpose": "зачем этот слайд нужен",
      "visual_hint": "идея визуала/схемы/иконок/стрелок",
      "icon": "target|people|data|risk|idea|growth|shield|map|clock|check",
      "takeaway": "главный вывод слайда",
      "notes": "Короткая заметка докладчика"
    }}
  ],
  "conclusion": "Ключевой вывод презентации"
}}"""


async def _generate_structure(
    topic: str,
    slides_count: int,
    language: str,
) -> PresentationData:
    """Вызывает LLM для генерации структуры презентации."""
    settings = get_settings()
    lang_name = "русском" if language == "ru" else "English"
    prompt = _STRUCTURE_PROMPT.format(
        topic=topic,
        slides_count=slides_count,
        language=lang_name,
    )

    logger.info("PPTX: generating structure for topic=%r slides=%d", topic[:60], slides_count)

    try:
        async with httpx.AsyncClient(timeout=httpx.Timeout(30.0, read=60.0)) as client:
            response = await client.post(
                f"{settings.MWS_BASE_URL}/chat/completions",
                headers={
                    "Authorization": f"Bearer {settings.MWS_API_KEY}",
                    "Content-Type": "application/json",
                },
                json={
                    "model": settings.LLM_MODEL,  # Топовая модель для качества
                    "messages": [{"role": "user", "content": prompt}],
                    "temperature": 0.6,
                    "max_tokens": 2048,
                    "stream": False,
                },
            )
            response.raise_for_status()
            data = response.json()
    except httpx.TimeoutException as e:
        raise ValueError(f"LLM timeout при генерации структуры: {e}") from e
    except Exception as e:
        raise ValueError(f"LLM API error: {e}") from e

    raw = data["choices"][0]["message"]["content"].strip()

    # Убираем <think>...</think> и code fences
    raw = re.sub(r"<think>[\s\S]*?</think>", "", raw).strip()
    raw = re.sub(r"```(?:json)?\s*", "", raw).strip("` \n")

    # Извлекаем JSON
    start = raw.find("{")
    end = raw.rfind("}")
    if start == -1 or end == -1:
        raise ValueError(f"LLM не вернул JSON структуру. Ответ: {raw[:200]}")

    try:
        obj = json.loads(raw[start : end + 1])
    except json.JSONDecodeError as e:
        raise ValueError(f"Невалидный JSON от LLM: {e}. Raw: {raw[:300]}") from e

    # Normalize slides
    slides = []
    for s in obj.get("slides", []):
        slides.append(
            SlideData(
                title=str(s.get("title", "Слайд")),
                bullets=[str(b) for b in s.get("bullets", [])],
                notes=str(s.get("notes", "")),
                purpose=str(s.get("purpose", "")),
                visual_hint=str(s.get("visual_hint", "")),
                layout=str(s.get("layout", "")),
                highlight=str(s.get("highlight", "")),
                body=str(s.get("body", "")),
                evidence=[
                    str(e).strip()
                    for e in s.get("evidence", [])
                    if str(e).strip()
                ][:4]
                if isinstance(s.get("evidence", []), list)
                else [],
                takeaway=str(s.get("takeaway", "")),
                icon=str(s.get("icon", "")),
            )
        )

    logger.info("PPTX: LLM generated %d slides", len(slides))

    return PresentationData(
        presentation_title=str(obj.get("presentation_title", topic)),
        subtitle=str(obj.get("subtitle", "")),
        slides=slides,
        conclusion=str(obj.get("conclusion", "")),
    )


# ──────────────────────────────────────────────────────────────────
# Создание PPTX-файла
# ──────────────────────────────────────────────────────────────────


def _build_pptx(data: PresentationData, style: str, file_path: Path) -> None:
    """
    Создаёт PPTX-файл с корпоративным дизайном МТС.

    Цвета:
      MTS Red:   #E30613
      MTS Dark:  #231F20
      MTS White: #FFFFFF
      MTS Gray:  #F5F5F5
    """
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches, Pt
    except ImportError as e:
        raise RuntimeError(
            "python-pptx не установлен. Добавьте 'python-pptx' в requirements.txt."
        ) from e

    # ── Цветовая палитра ──
    style = (style or "modern_light").strip()
    if style == "dark":
        BG_COLOR = RGBColor(0x1C, 0x1C, 0x1C)
        TITLE_COLOR = RGBColor(0xFF, 0xFF, 0xFF)
        BODY_COLOR = RGBColor(0xE0, 0xE0, 0xE0)
        MUTED_COLOR = RGBColor(0x9A, 0x9A, 0x9A)
        PANEL_COLOR = RGBColor(0x2A, 0x2A, 0x2A)
        ACCENT_COLOR = RGBColor(0xE3, 0x06, 0x13)
    elif style == "minimal":
        BG_COLOR = RGBColor(0xFF, 0xFF, 0xFF)
        TITLE_COLOR = RGBColor(0x23, 0x1F, 0x20)
        BODY_COLOR = RGBColor(0x44, 0x44, 0x44)
        MUTED_COLOR = RGBColor(0x88, 0x88, 0x88)
        PANEL_COLOR = RGBColor(0xF3, 0xF3, 0xF3)
        ACCENT_COLOR = RGBColor(0xE3, 0x06, 0x13)
    elif style == "strict_corporate":
        BG_COLOR = RGBColor(0xFF, 0xFF, 0xFF)
        TITLE_COLOR = RGBColor(0x1C, 0x1C, 0x1C)
        BODY_COLOR = RGBColor(0x2E, 0x2E, 0x2E)
        MUTED_COLOR = RGBColor(0x70, 0x70, 0x70)
        PANEL_COLOR = RGBColor(0xF4, 0xF5, 0xF7)
        ACCENT_COLOR = RGBColor(0xE3, 0x06, 0x13)
    elif style == "mts_corporate":
        BG_COLOR = RGBColor(0xF5, 0xF5, 0xF5)
        TITLE_COLOR = RGBColor(0x23, 0x1F, 0x20)
        BODY_COLOR = RGBColor(0x33, 0x33, 0x33)
        MUTED_COLOR = RGBColor(0x70, 0x70, 0x70)
        PANEL_COLOR = RGBColor(0xFF, 0xFF, 0xFF)
        ACCENT_COLOR = RGBColor(0xE3, 0x06, 0x13)
    elif style == "creative":
        BG_COLOR = RGBColor(0xFF, 0xF4, 0xF4)
        TITLE_COLOR = RGBColor(0xE3, 0x06, 0x13)
        BODY_COLOR = RGBColor(0x23, 0x1F, 0x20)
        MUTED_COLOR = RGBColor(0x7A, 0x38, 0x38)
        PANEL_COLOR = RGBColor(0xFF, 0xFF, 0xFF)
        ACCENT_COLOR = RGBColor(0xE3, 0x06, 0x13)
    elif style == "tech":
        BG_COLOR = RGBColor(0xF7, 0xF9, 0xFC)
        TITLE_COLOR = RGBColor(0x18, 0x24, 0x35)
        BODY_COLOR = RGBColor(0x2C, 0x3E, 0x50)
        MUTED_COLOR = RGBColor(0x65, 0x74, 0x86)
        PANEL_COLOR = RGBColor(0xFF, 0xFF, 0xFF)
        ACCENT_COLOR = RGBColor(0x00, 0x72, 0xCE)
    elif style == "editorial":
        BG_COLOR = RGBColor(0xFB, 0xFA, 0xF8)
        TITLE_COLOR = RGBColor(0x18, 0x18, 0x18)
        BODY_COLOR = RGBColor(0x34, 0x34, 0x34)
        MUTED_COLOR = RGBColor(0x7A, 0x72, 0x68)
        PANEL_COLOR = RGBColor(0xFF, 0xFF, 0xFF)
        ACCENT_COLOR = RGBColor(0xB0, 0x1E, 0x2F)
    elif style == "modern_light":
        BG_COLOR = RGBColor(0xFA, 0xFB, 0xFC)
        TITLE_COLOR = RGBColor(0x20, 0x24, 0x2A)
        BODY_COLOR = RGBColor(0x3C, 0x44, 0x4F)
        MUTED_COLOR = RGBColor(0x78, 0x84, 0x91)
        PANEL_COLOR = RGBColor(0xFF, 0xFF, 0xFF)
        ACCENT_COLOR = RGBColor(0xE3, 0x06, 0x13)
    else:  # corporate (default)
        BG_COLOR = RGBColor(0xFA, 0xFB, 0xFC)
        TITLE_COLOR = RGBColor(0x23, 0x1F, 0x20)
        BODY_COLOR = RGBColor(0x33, 0x33, 0x33)
        MUTED_COLOR = RGBColor(0x78, 0x84, 0x91)
        PANEL_COLOR = RGBColor(0xFF, 0xFF, 0xFF)
        ACCENT_COLOR = RGBColor(0xE3, 0x06, 0x13)

    RED = RGBColor(0xE3, 0x06, 0x13)  # MTS Red всегда
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    BLACK = RGBColor(0x18, 0x18, 0x18)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # Пустой макет

    def _set_bg(slide, color: RGBColor):
        """Устанавливает фон слайда."""
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _font_size_for(text: str, base: int, limit: int, minimum: int = 18) -> int:
        overflow = max(0, len(text or "") - limit)
        return max(minimum, base - overflow // 8)

    def _box_capacity(width, height, font_size: int) -> int:
        return max(28, int(width.inches * height.inches * 15 * (18 / max(font_size, 8))))

    def _add_text_box(slide, text: str, left, top, width, height, font_size: int,
                      bold: bool = False, color: RGBColor = None, align=None,
                      valign=None):
        """Добавляет текстовый блок."""
        text = _shorten_text(text, _box_capacity(width, height, font_size))
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.06)
        tf.margin_right = Inches(0.06)
        tf.margin_top = Inches(0.03)
        tf.margin_bottom = Inches(0.03)
        if valign:
            tf.vertical_anchor = valign
        p = tf.paragraphs[0]
        p.text = text
        if align:
            p.alignment = align
        run = p.runs[0]
        run.font.size = Pt(font_size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
        return txBox

    def _add_shape(slide, shape_type, left, top, width, height, fill, line=None, radius=False):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE if radius else shape_type,
            left,
            top,
            width,
            height,
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        if line:
            shape.line.color.rgb = line
            shape.line.width = Pt(1)
        else:
            shape.line.fill.background()
        return shape

    def _add_card(slide, left, top, width, height, fill=None):
        return _add_shape(
            slide,
            MSO_SHAPE.RECTANGLE,
            left,
            top,
            width,
            height,
            fill or PANEL_COLOR,
            line=RGBColor(0xE5, 0xE7, 0xEB),
            radius=True,
        )

    def _add_icon(slide, icon: str, left, top, size, color=None):
        glyphs = {
            "target": "◎",
            "people": "◇",
            "data": "▦",
            "risk": "!",
            "idea": "✦",
            "growth": "↗",
            "shield": "◆",
            "map": "⌁",
            "clock": "◷",
            "check": "✓",
        }
        text = glyphs.get((icon or "").lower(), "•")
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
        circle.fill.solid()
        circle.fill.fore_color.rgb = color or ACCENT_COLOR
        circle.line.fill.background()
        return _add_text_box(
            slide,
            text,
            left,
            top + Inches(0.02),
            size,
            size,
            font_size=18,
            bold=True,
            color=WHITE,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )

    def _add_header(slide, title: str, idx: int):
        _add_text_box(
            slide,
            title,
            Inches(0.65),
            Inches(0.35),
            Inches(10.7),
            Inches(0.75),
            font_size=_font_size_for(title, 30, 52, 20),
            bold=True,
            color=TITLE_COLOR,
        )
        _add_shape(
            slide,
            MSO_SHAPE.RECTANGLE,
            Inches(0.65),
            Inches(1.18),
            Inches(2.0),
            Inches(0.04),
            ACCENT_COLOR,
        )
        _add_text_box(
            slide,
            f"{idx + 1:02}",
            Inches(12.25),
            Inches(0.38),
            Inches(0.55),
            Inches(0.35),
            font_size=12,
            bold=True,
            color=ACCENT_COLOR,
            align=PP_ALIGN.RIGHT,
        )

    def _split_bullets(bullets: list[str]) -> tuple[str, list[str]]:
        clean = [str(b).strip() for b in bullets if str(b).strip()]
        if not clean:
            return "", []
        return clean[0], clean[1:5]

    def _first_text(*values) -> str:
        for value in values:
            if isinstance(value, list):
                for item in value:
                    if str(item).strip():
                        return str(item).strip()
            elif str(value or "").strip():
                return str(value).strip()
        return ""

    def _add_wrapped_list(slide, items: list[str], left, top, width, font_size=14, color=None, prefix="•"):
        y = top
        for item in [str(i).strip() for i in items if str(i).strip()][:4]:
            _add_text_box(slide, f"{prefix} {item}", left, y, width, Inches(0.48), font_size, color=color or BODY_COLOR)
            y += Inches(0.55)

    def _layout_hero(slide, slide_data, idx: int):
        _add_header(slide, slide_data.title, idx)
        lead, rest = _split_bullets(slide_data.bullets)
        highlight = slide_data.highlight or lead or slide_data.purpose
        _add_text_box(
            slide,
            highlight,
            Inches(0.9),
            Inches(1.75),
            Inches(7.1),
            Inches(1.8),
            font_size=_font_size_for(highlight, 31, 70, 21),
            bold=True,
            color=TITLE_COLOR,
            valign=MSO_ANCHOR.MIDDLE,
        )
        _add_card(slide, Inches(8.45), Inches(1.65), Inches(3.9), Inches(3.7))
        _add_icon(slide, slide_data.icon, Inches(8.75), Inches(1.95), Inches(0.55))
        body = _first_text(slide_data.body, slide_data.evidence, rest, slide_data.purpose)
        _add_text_box(slide, body, Inches(8.85), Inches(2.65), Inches(3.05), Inches(1.25), 14, color=BODY_COLOR)
        _add_wrapped_list(slide, rest[:2] or slide_data.evidence[:2], Inches(8.85), Inches(4.08), Inches(3.05), 12, MUTED_COLOR, "→")
        if slide_data.takeaway:
            _add_text_box(
                slide,
                slide_data.takeaway,
                Inches(0.95),
                Inches(5.45),
                Inches(10.9),
                Inches(0.75),
                18,
                bold=True,
                color=ACCENT_COLOR,
                align=PP_ALIGN.CENTER,
            )

    def _layout_split(slide, slide_data, idx: int):
        _add_header(slide, slide_data.title, idx)
        highlight = slide_data.highlight or slide_data.purpose or (slide_data.bullets[0] if slide_data.bullets else "")
        _add_card(slide, Inches(0.8), Inches(1.55), Inches(5.1), Inches(4.65))
        _add_text_box(
            slide,
            highlight,
            Inches(1.15),
            Inches(2.05),
            Inches(4.35),
            Inches(2.1),
            _font_size_for(highlight, 25, 62, 18),
            bold=True,
            color=TITLE_COLOR,
            valign=MSO_ANCHOR.MIDDLE,
        )
        if slide_data.takeaway:
            _add_text_box(slide, slide_data.takeaway, Inches(1.15), Inches(4.65), Inches(4.35), Inches(0.75), 15, color=MUTED_COLOR)
        if slide_data.body:
            _add_text_box(slide, slide_data.body, Inches(6.55), Inches(1.55), Inches(5.15), Inches(0.85), 14, color=MUTED_COLOR)
        y = 2.58 if slide_data.body else 1.75
        for item in slide_data.bullets[:3]:
            _add_icon(slide, "check", Inches(6.55), Inches(y + 0.02), Inches(0.32), ACCENT_COLOR)
            _add_text_box(slide, item, Inches(7.05), Inches(y), Inches(4.8), Inches(0.55), 16, color=BODY_COLOR)
            y += 0.9
        if slide_data.evidence:
            _add_text_box(slide, "Факт: " + slide_data.evidence[0], Inches(7.05), Inches(5.85), Inches(4.8), Inches(0.42), 12, color=MUTED_COLOR)

    def _layout_cards(slide, slide_data, idx: int):
        _add_header(slide, slide_data.title, idx)
        if slide_data.highlight:
            _add_text_box(slide, slide_data.highlight, Inches(0.8), Inches(1.35), Inches(11.5), Inches(0.55), 17, bold=True, color=ACCENT_COLOR)
        bullets = slide_data.bullets[:4] or [slide_data.purpose]
        card_w = 5.55 if len(bullets) <= 2 else 2.75
        for n, item in enumerate(bullets):
            col = n % 4 if len(bullets) > 2 else n
            row = 0 if len(bullets) > 2 else n // 2
            left = Inches(0.75 + col * 3.05) if len(bullets) > 2 else Inches(0.8 + col * 5.9)
            top = Inches(2.0 + row * 1.95)
            _add_card(slide, left, top, Inches(card_w), Inches(2.55 if len(bullets) <= 2 else 3.25))
            _add_icon(slide, slide_data.icon, left + Inches(0.25), top + Inches(0.25), Inches(0.42))
            _add_text_box(slide, item, left + Inches(0.28), top + Inches(0.95), Inches(card_w - 0.45), Inches(1.25), 15, bold=True, color=BODY_COLOR)
            evidence = slide_data.evidence[n] if n < len(slide_data.evidence) else ""
            if evidence:
                _add_text_box(slide, evidence, left + Inches(0.28), top + Inches(2.05), Inches(card_w - 0.45), Inches(0.65), 11, color=MUTED_COLOR)
        if slide_data.takeaway:
            _add_text_box(slide, slide_data.takeaway, Inches(0.9), Inches(6.25), Inches(11.4), Inches(0.45), 14, color=MUTED_COLOR, align=PP_ALIGN.CENTER)

    def _layout_process(slide, slide_data, idx: int):
        _add_header(slide, slide_data.title, idx)
        steps = slide_data.bullets[:4] or [slide_data.purpose]
        y = Inches(3.05)
        for n, item in enumerate(steps):
            left = Inches(0.75 + n * 3.05)
            _add_shape(slide, MSO_SHAPE.OVAL, left, y, Inches(0.75), Inches(0.75), ACCENT_COLOR)
            _add_text_box(slide, str(n + 1), left, y + Inches(0.07), Inches(0.75), Inches(0.45), 18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            _add_text_box(slide, item, left - Inches(0.15), y + Inches(1.0), Inches(2.3), Inches(1.0), 14, bold=True, color=BODY_COLOR, align=PP_ALIGN.CENTER)
            if n < len(steps) - 1:
                _add_text_box(slide, "→", left + Inches(1.15), y + Inches(0.02), Inches(1.05), Inches(0.7), 30, bold=True, color=ACCENT_COLOR, align=PP_ALIGN.CENTER)
        if slide_data.highlight:
            _add_text_box(slide, slide_data.highlight, Inches(1.1), Inches(1.55), Inches(10.6), Inches(0.8), 22, bold=True, color=TITLE_COLOR, align=PP_ALIGN.CENTER)
        if slide_data.body:
            _add_text_box(slide, slide_data.body, Inches(1.1), Inches(5.75), Inches(10.6), Inches(0.55), 13, color=MUTED_COLOR, align=PP_ALIGN.CENTER)

    def _layout_comparison(slide, slide_data, idx: int):
        _add_header(slide, slide_data.title, idx)
        bullets = slide_data.bullets[:4]
        left_items = bullets[:2] or [slide_data.purpose]
        right_items = bullets[2:4] or ([slide_data.takeaway] if slide_data.takeaway else [])
        _add_card(slide, Inches(0.85), Inches(1.75), Inches(5.45), Inches(4.55))
        _add_card(slide, Inches(7.0), Inches(1.75), Inches(5.45), Inches(4.55))
        _add_text_box(slide, slide_data.highlight or "Сейчас", Inches(1.15), Inches(2.05), Inches(4.8), Inches(0.55), 19, bold=True, color=TITLE_COLOR)
        _add_text_box(slide, "Что меняется", Inches(7.3), Inches(2.05), Inches(4.8), Inches(0.55), 19, bold=True, color=ACCENT_COLOR)
        if slide_data.body:
            _add_text_box(slide, slide_data.body, Inches(1.2), Inches(5.35), Inches(10.8), Inches(0.55), 12, color=MUTED_COLOR, align=PP_ALIGN.CENTER)
        y = 2.85
        for item in left_items:
            _add_text_box(slide, f"• {item}", Inches(1.2), Inches(y), Inches(4.7), Inches(0.65), 15, color=BODY_COLOR)
            y += 0.75
        y = 2.85
        for item in right_items:
            _add_text_box(slide, f"→ {item}", Inches(7.35), Inches(y), Inches(4.7), Inches(0.65), 15, color=BODY_COLOR)
            y += 0.75

    def _layout_quote(slide, slide_data, idx: int):
        _add_header(slide, slide_data.title, idx)
        quote = slide_data.highlight or slide_data.takeaway or (slide_data.bullets[0] if slide_data.bullets else slide_data.purpose)
        _add_text_box(slide, "“", Inches(0.9), Inches(1.4), Inches(0.8), Inches(0.8), 50, bold=True, color=ACCENT_COLOR)
        _add_text_box(slide, quote, Inches(1.5), Inches(1.85), Inches(9.9), Inches(2.0), _font_size_for(quote, 30, 80, 20), bold=True, color=TITLE_COLOR, align=PP_ALIGN.CENTER)
        y = 4.55
        quote_items = slide_data.evidence[:2] + slide_data.bullets[:2]
        for item in quote_items[:3]:
            _add_text_box(slide, item, Inches(2.25), Inches(y), Inches(8.8), Inches(0.45), 15, color=BODY_COLOR, align=PP_ALIGN.CENTER)
            y += 0.55

    def _layout_takeaway(slide, slide_data, idx: int):
        _add_header(slide, slide_data.title, idx)
        _add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.0), Inches(5.15), prs.slide_width, Inches(2.35), ACCENT_COLOR)
        msg = slide_data.takeaway or slide_data.highlight or (slide_data.bullets[0] if slide_data.bullets else slide_data.purpose)
        _add_text_box(slide, msg, Inches(1.0), Inches(5.45), Inches(11.2), Inches(1.0), _font_size_for(msg, 27, 80, 18), bold=True, color=WHITE, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        y = 1.65
        takeaway_items = ([slide_data.body] if slide_data.body else []) + slide_data.bullets[:3]
        for item in takeaway_items[:3]:
            _add_icon(slide, "check", Inches(1.4), Inches(y + 0.02), Inches(0.35), ACCENT_COLOR)
            _add_text_box(slide, item, Inches(1.95), Inches(y), Inches(9.7), Inches(0.5), 17, color=BODY_COLOR)
            y += 0.82

    def _conclusion_points() -> list[str]:
        raw_parts = [
            p.strip(" -•\n\t")
            for p in re.split(r"(?<=[.!?])\s+|\n+", data.conclusion or "")
            if len(p.strip(" -•\n\t")) > 6
        ]
        slide_takeaways = [
            s.takeaway.strip()
            for s in data.slides
            if getattr(s, "takeaway", "") and len(s.takeaway.strip()) > 6
        ]
        slide_highlights = [
            s.highlight.strip()
            for s in data.slides
            if getattr(s, "highlight", "") and len(s.highlight.strip()) > 6
        ]
        points = raw_parts + slide_takeaways + slide_highlights
        result = []
        for point in points:
            if point not in result:
                result.append(point)
            if len(result) >= 3:
                break
        return result or [
            "Материал собран в единую историю с понятной логикой.",
            "Ключевые факты связаны с целью и аудиторией презентации.",
            "Финальный вывод показывает, что нужно запомнить после выступления.",
        ]

    # ── Слайд 0: Титульный ──
    title_slide = prs.slides.add_slide(blank_layout)
    title_bg = RED if style == "mts_corporate" else ACCENT_COLOR
    _set_bg(title_slide, title_bg)

    # Белый прямоугольник снизу (30% высоты)
    bar = title_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(5.2),
        prs.slide_width, Inches(2.3),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = WHITE
    bar.line.fill.background()

    # Заголовок на красном фоне
    _add_text_box(
        title_slide, data.presentation_title,
        Inches(0.8), Inches(1.2),
        Inches(11.7), Inches(2.5),
        font_size=_font_size_for(data.presentation_title, 36, 58, 22), bold=True, color=WHITE,
        align=PP_ALIGN.LEFT,
    )

    if data.subtitle:
        _add_text_box(
            title_slide, data.subtitle,
            Inches(0.8), Inches(3.8),
            Inches(11.7), Inches(1.0),
            font_size=_font_size_for(data.subtitle, 18, 82, 13), bold=False, color=WHITE,
            align=PP_ALIGN.LEFT,
        )

    # Надпись в белой полосе
    brand_label = "МТС" if style == "mts_corporate" else "Презентация"
    _add_text_box(
        title_slide, brand_label,
        Inches(0.8), Inches(5.4),
        Inches(5), Inches(0.8),
        font_size=24, bold=True, color=ACCENT_COLOR,
        align=PP_ALIGN.LEFT,
    )

    # ── Слайды контента ──
    for i, slide_data in enumerate(data.slides):
        slide = prs.slides.add_slide(blank_layout)
        _set_bg(slide, BG_COLOR)

        if style == "mts_corporate":
            _add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.12), prs.slide_height, RED)

        layout = (slide_data.layout or "").lower().strip()
        if not layout:
            layout = ["hero", "split", "cards", "process", "comparison", "quote", "takeaway"][i % 7]

        if layout == "hero":
            _layout_hero(slide, slide_data, i)
        elif layout in {"split", "two_column", "two-column"}:
            _layout_split(slide, slide_data, i)
        elif layout in {"cards", "card_grid"}:
            _layout_cards(slide, slide_data, i)
        elif layout in {"timeline", "process", "flow"}:
            _layout_process(slide, slide_data, i)
        elif layout == "comparison":
            _layout_comparison(slide, slide_data, i)
        elif layout == "quote":
            _layout_quote(slide, slide_data, i)
        elif layout == "takeaway":
            _layout_takeaway(slide, slide_data, i)
        else:
            _layout_split(slide, slide_data, i)

        # Заметки докладчика
        if slide_data.notes:
            notes_slide = slide.notes_slide
            tf = notes_slide.notes_text_frame
            tf.text = slide_data.notes

    # ── Финальный слайд: Заключение ──
    final_slide = prs.slides.add_slide(blank_layout)
    _set_bg(final_slide, BG_COLOR)

    _add_text_box(
        final_slide, "Главный вывод",
        Inches(0.75), Inches(0.45),
        Inches(11.7), Inches(0.7),
        font_size=30, bold=True, color=TITLE_COLOR,
    )
    _add_shape(
        final_slide,
        MSO_SHAPE.RECTANGLE,
        Inches(0.75),
        Inches(1.23),
        Inches(2.0),
        Inches(0.04),
        ACCENT_COLOR,
    )

    for idx, point in enumerate(_conclusion_points(), 1):
        top = Inches(1.65 + (idx - 1) * 1.55)
        _add_card(final_slide, Inches(0.95), top, Inches(11.25), Inches(1.2))
        _add_icon(final_slide, "check", Inches(1.25), top + Inches(0.28), Inches(0.45), ACCENT_COLOR)
        _add_text_box(
            final_slide,
            point,
            Inches(1.95),
            top + Inches(0.22),
            Inches(9.75),
            Inches(0.68),
            _font_size_for(point, 18, 110, 13),
            bold=True,
            color=BODY_COLOR,
            valign=MSO_ANCHOR.MIDDLE,
        )

    _add_text_box(final_slide, "Спасибо", Inches(0.95), Inches(6.45), Inches(11.2), Inches(0.35), 14, color=MUTED_COLOR, align=PP_ALIGN.CENTER)

    prs.save(str(file_path))
    logger.info("PPTX: saved to %s (%d bytes)", file_path, file_path.stat().st_size)


def _preview_markdown(pptx_data: PresentationData) -> str:
    """Собирает markdown-превью структуры для ответа в чат."""
    preview_lines = [
        f"## 📊 {pptx_data.presentation_title}",
        f"*{pptx_data.subtitle}*" if pptx_data.subtitle else "",
        "",
        "### Содержание:",
    ]
    for i, slide in enumerate(pptx_data.slides, 1):
        preview_lines.append(f"**{i}. {slide.title}**")
        for bullet in slide.bullets[:2]:
            preview_lines.append(f"  - {bullet}")

    if pptx_data.conclusion:
        preview_lines += ["", f"### Вывод:\n> {pptx_data.conclusion}"]

    return "\n".join(filter(None, preview_lines))


async def _render_and_register_pptx(
    pptx_data: PresentationData,
    *,
    topic: str,
    style: str,
    user_id: str,
) -> JSONResponse:
    """Рендерит PPTX, регистрирует файл и возвращает OpenAPI-ответ."""
    file_id = uuid.uuid4().hex[:16]
    safe_title = (
        re.sub(r"[^\w\s-]", "", pptx_data.presentation_title)[:50]
        .strip()
        .replace(" ", "_")
    )
    if not safe_title:
        safe_title = "presentation"
    filename = f"{safe_title}_{file_id}.pptx"
    file_path = PPTX_DIR / filename

    try:
        await asyncio.to_thread(_build_pptx, pptx_data, style, file_path)
    except RuntimeError as e:
        logger.error("PPTX build failed: %s", e)
        return JSONResponse(status_code=500, content={"error": str(e)})
    except Exception:
        logger.exception("PPTX build unexpected error")
        return JSONResponse(
            status_code=500,
            content={"error": "Внутренняя ошибка создания файла."},
        )

    total_slides = len(pptx_data.slides) + 2
    _file_registry[file_id] = {
        "path": str(file_path),
        "filename": filename,
        "user_id": user_id,
        "topic": topic,
        "created_at": time.time(),
        "slide_count": total_slides,
    }

    download_url = f"/v1/presentations/download/{file_id}"

    logger.info(
        "PPTX ready: file_id=%s slides=%d size=%dKB",
        file_id,
        total_slides,
        file_path.stat().st_size // 1024,
    )

    return JSONResponse(
        content={
            "file_id": file_id,
            "download_url": download_url,
            "filename": filename,
            "slide_count": total_slides,
            "topic": topic,
            "preview_markdown": _preview_markdown(pptx_data),
        }
    )


def _outline_to_presentation_data(
    request: PresentationFromOutlineRequest,
) -> PresentationData:
    """Преобразует согласованный outline агента в формат python-pptx."""
    brief = request.brief or {}
    topic = _compact_title(str(brief.get("deck_title") or brief.get("topic") or "Презентация"), 72)
    audience = str(brief.get("audience") or "").strip()
    goal = str(brief.get("goal") or "").strip()
    subtitle = _compact_subtitle(audience, goal, str(brief.get("deck_subtitle") or ""))

    slides: list[SlideData] = []
    for idx, item in enumerate(request.outline, 1):
        bullets = [str(b).strip() for b in item.bullets if str(b).strip()]
        if item.purpose and item.purpose not in bullets:
            bullets.insert(0, item.purpose)
        if not bullets:
            bullets = ["Ключевая мысль", "Аргумент", "Практический вывод"]
        bullets = [_shorten_text(b, 88) for b in bullets][:4]

        notes_parts = [item.speaker_notes.strip()]
        if item.visual_hint:
            notes_parts.append(f"Визуальная идея: {item.visual_hint}")
        notes = "\n".join(part for part in notes_parts if part)

        slides.append(
            SlideData(
                title=item.title.strip() or f"Слайд {idx}",
                bullets=bullets[:5],
                notes=notes,
                purpose=_shorten_text(item.purpose, 110),
                visual_hint=_shorten_text(item.visual_hint, 120),
                layout=item.layout.strip(),
                highlight=_shorten_text(item.highlight, 80),
                body=_shorten_text(item.body, 170),
                evidence=[_shorten_text(e, 90) for e in item.evidence if str(e).strip()][:3],
                takeaway=_shorten_text(item.takeaway, 95),
                icon=item.icon.strip(),
            )
        )

    closing_message = str(brief.get("closing_message") or "").strip()
    if closing_message and len(closing_message) > 12:
        conclusion = closing_message
    else:
        takeaways = [
            slide.takeaway
            for slide in slides
            if slide.takeaway and len(slide.takeaway.strip()) > 8
        ][:3]
        if takeaways:
            conclusion = " ".join(takeaways)
        elif goal and len(goal) > 12:
            conclusion = f"Презентация подводит аудиторию к выводу: {goal}."
        else:
            conclusion = "Ключевые факты собраны в единую логику, а выводы показывают, что аудитории нужно запомнить и сделать дальше."

    return PresentationData(
        presentation_title=topic,
        subtitle=subtitle,
        slides=slides,
        conclusion=conclusion,
    )


# ──────────────────────────────────────────────────────────────────
# Роуты
# ──────────────────────────────────────────────────────────────────


@router.post("/presentations/generate")
async def generate_presentation(request: PresentationRequest):
    """
    Генерирует профессиональную PPTX-презентацию по заданной теме.

    Workflow:
      1. Вызов LLM для создания структуры (заголовки, буллиты, выводы)
      2. Рендеринг PPTX с корпоративным дизайном МТС
      3. Возврат ссылки на скачивание

    Returns:
        {
            "file_id": "...",
            "download_url": "/v1/presentations/download/{file_id}",
            "filename": "...",
            "slide_count": N,
            "topic": "...",
            "preview_markdown": "# ...\n..."
        }
    """
    logger.info(
        "presentations/generate: topic=%r slides=%d style=%s user=%s",
        request.topic[:60], request.slides_count, request.style, request.user_id,
    )

    # ── Генерация структуры через LLM ──
    try:
        pptx_data = await _generate_structure(
            request.topic, request.slides_count, request.language
        )
    except ValueError as e:
        logger.error("PPTX structure generation failed: %s", e)
        return JSONResponse(
            status_code=502,
            content={"error": f"Ошибка генерации структуры: {str(e)[:300]}"},
        )

    return await _render_and_register_pptx(
        pptx_data,
        topic=request.topic,
        style=request.style,
        user_id=request.user_id,
    )


@router.post("/presentations/generate-from-outline")
async def generate_presentation_from_outline(request: PresentationFromOutlineRequest):
    """
    Генерирует PPTX из уже согласованного брифа и структуры слайдов.

    Используется интерактивным presentation-agent после подтверждения
    структуры пользователем.
    """
    logger.info(
        "presentations/generate-from-outline: slides=%d style=%s user=%s",
        len(request.outline),
        request.style,
        request.user_id,
    )

    if not request.outline:
        return JSONResponse(
            status_code=400,
            content={"error": "Нет структуры слайдов для генерации PPTX."},
        )

    pptx_data = _outline_to_presentation_data(request)
    topic = str(request.brief.get("topic") or pptx_data.presentation_title)

    return await _render_and_register_pptx(
        pptx_data,
        topic=topic,
        style=request.style,
        user_id=request.user_id,
    )


@router.get("/presentations/download/{file_id}")
async def download_presentation(file_id: str):
    """Скачивает готовый PPTX-файл по ID."""
    info = _file_registry.get(file_id)
    if not info:
        return JSONResponse(status_code=404, content={"error": "Файл не найден."})

    path = Path(info["path"])
    if not path.exists():
        return JSONResponse(
            status_code=410,
            content={"error": "Файл был удалён с сервера."},
        )

    return FileResponse(
        path=str(path),
        filename=info["filename"],
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )


@router.get("/presentations/list/{user_id}")
async def list_presentations(user_id: str):
    """Возвращает список презентаций пользователя."""
    user_files = [
        {
            "file_id": fid,
            "filename": info["filename"],
            "topic": info["topic"],
            "slide_count": info["slide_count"],
            "download_url": f"/v1/presentations/download/{fid}",
            "created_at": info["created_at"],
        }
        for fid, info in _file_registry.items()
        if info.get("user_id") == user_id
    ]
    return JSONResponse(content={"presentations": user_files, "count": len(user_files)})
